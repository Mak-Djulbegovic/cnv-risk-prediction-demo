"""
AI for Ophthalmology Research - Clean Clinical Rebuild
======================================================
Native PowerPoint build. No embedded diagram screenshots, no emoji glyphs.
Editorial / "journal" design language: white slides, navy headers, one teal
accent, generous whitespace. Statistical figures are the clean title-less
versions (figure*_clean.png).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = '/Users/makdjulbegovic/Desktop/Claude_Teach'

# ---- Palette --------------------------------------------------------------
NAVY    = RGBColor(0x16, 0x36, 0x5C)
TEAL    = RGBColor(0x2C, 0x7A, 0x7B)
CRIMSON = RGBColor(0xC0, 0x39, 0x2B)
AMBER   = RGBColor(0xB9, 0x77, 0x0E)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
INK     = RGBColor(0x21, 0x25, 0x29)
GRAY    = RGBColor(0x5A, 0x65, 0x70)
MUTE    = RGBColor(0x8A, 0x93, 0x9E)
RULE    = RGBColor(0xDD, 0xE2, 0xE7)
LIGHT   = RGBColor(0xF5, 0xF7, 0xF9)
PANEL   = RGBColor(0x1E, 0x2A, 0x38)   # dark code panel
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = 'Helvetica Neue'
BODY = 'Helvetica Neue'
MONO = 'Menlo'

# ---- Canvas (16:9) --------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)
CONTENT_W = SW - 2 * MARGIN
BLANK = prs.slide_layouts[6]


# ---- Low-level helpers ----------------------------------------------------
def _set_font(run, size, color, bold=False, italic=False, font=BODY):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tb, tf


def para(tf, text, size, color, bold=False, italic=False, font=BODY,
         align=PP_ALIGN.LEFT, space_before=0, space_after=6, first=False,
         line=None, link=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    run = p.add_run()
    run.text = text
    _set_font(run, size, color, bold, italic, font)
    if link:
        run.hyperlink.address = link
        _set_font(run, size, color, bold, italic, font)  # keep our color, not theme
    return p


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None,
         shape=MSO_SHAPE.RECTANGLE, shadow_off=True):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    if shadow_off:
        sp.shadow.inherit = False
    return sp


def gradient_bg(slide, c1, c2, angle=90):
    sp = rect(slide, 0, 0, SW, SH)
    sp.fill.gradient()
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    sp.fill.gradient_stops[0].color.rgb = c1
    sp.fill.gradient_stops[1].color.rgb = c2
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def shape_text(sp, lines, anchor=MSO_ANCHOR.MIDDLE):
    """lines: list of dicts {text,size,color,bold,italic,font,align,space_after}"""
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', PP_ALIGN.CENTER)
        p.space_before = Pt(ln.get('space_before', 0))
        p.space_after = Pt(ln.get('space_after', 4))
        r = p.add_run()
        r.text = ln['text']
        _set_font(r, ln['size'], ln['color'], ln.get('bold', False),
                  ln.get('italic', False), ln.get('font', BODY))


def add_slide():
    return prs.slides.add_slide(BLANK)


def header(slide, title, accent=TEAL, kicker=None):
    """Editorial header: optional kicker, title, short accent rule."""
    y = Inches(0.55)
    if kicker:
        _, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.3))
        para(tf, kicker.upper(), 13, accent, bold=True, font=HEAD, first=True,
             space_after=2)
        y = Inches(0.9)
    tb, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.7))
    para(tf, title, 30, NAVY, bold=True, font=HEAD, first=True, space_after=0)
    rule_y = y + Inches(0.62)
    rect(slide, MARGIN, rule_y, Inches(1.4), Pt(3), fill=accent)
    return Inches(1.85)  # content top


def footer(slide, num=None):
    # Auto-number from the slide's actual position so inserting/reordering
    # slides never desyncs the page numbers.
    num = len(prs.slides._sldIdLst)
    rect(slide, MARGIN, Inches(7.02), CONTENT_W, Pt(0.75), fill=RULE)
    _, tf = textbox(slide, MARGIN, Inches(7.08), Inches(8), Inches(0.3))
    para(tf, 'Synthetic data \u00b7 Educational demonstration only', 9, MUTE,
         italic=True, first=True, space_after=0)
    _, tf = textbox(slide, SW - MARGIN - Inches(1), Inches(7.08), Inches(1),
                    Inches(0.3), align=PP_ALIGN.RIGHT)
    para(tf, str(num), 10, MUTE, first=True, align=PP_ALIGN.RIGHT, space_after=0)


def card(slide, x, y, w, h, fill=LIGHT, line=RULE, line_w=Pt(1),
         radius=0.08):
    sp = rect(slide, x, y, w, h, fill=fill, line=line, line_w=line_w,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


# ---- helper: chat-style "prompt to Claude" callout ------------------------
# Shows the hypothetical natural-language prompt that produced the slide's
# artifact, so residents can see HOW each result was asked for.
def prompt_bar(slide, y, text, accent=TEAL, h=Inches(0.58)):
    card(slide, MARGIN, y, CONTENT_W, h, fill=LIGHT, line=RULE)
    rect(slide, MARGIN, y, Inches(0.12), h, fill=accent,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = textbox(slide, MARGIN + Inches(0.32), y, Inches(1.05), h,
                    anchor=MSO_ANCHOR.MIDDLE)
    para(tf, 'PROMPT', 11, accent, bold=True, font=HEAD, first=True,
         space_after=0)
    _, tf = textbox(slide, MARGIN + Inches(1.4), y, CONTENT_W - Inches(1.7), h,
                    anchor=MSO_ANCHOR.MIDDLE)
    para(tf, '\u201c' + text + '\u201d', 13, INK, italic=True, font=BODY,
         first=True, space_after=0, line=1.05)
    return y + h


# ===========================================================================
# SLIDE 1 - Title
# ===========================================================================
s = add_slide()
gradient_bg(s, NAVY, RGBColor(0x0C, 0x1E, 0x33))
# thin teal accent bar
rect(s, 0, Inches(3.02), SW, Pt(2.5), fill=TEAL)
_, tf = textbox(s, Inches(1.0), Inches(2.05), Inches(11.3), Inches(1.0))
para(tf, 'AI-POWERED RESEARCH IN OPHTHALMOLOGY', 15, RGBColor(0x9F, 0xC5, 0xC6),
     bold=True, font=HEAD, first=True, align=PP_ALIGN.CENTER, space_after=0)
_, tf = textbox(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.4))
para(tf, 'Agentic Systems: From Research Question to Clinical Tool',
     34, WHITE, bold=True, font=HEAD, first=True, align=PP_ALIGN.CENTER,
     space_after=0, line=1.05)
_, tf = textbox(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.5))
para(tf, 'A live demonstration with Claude Code', 18,
     RGBColor(0xC7, 0xD3, 0xE0), font=BODY, first=True, align=PP_ALIGN.CENTER,
     space_after=0)
# synthetic-data chip
chip = card(s, Inches(4.27), Inches(6.45), Inches(4.8), Inches(0.5),
            fill=RGBColor(0x24, 0x3B, 0x55), line=TEAL, line_w=Pt(1))
shape_text(chip, [{'text': 'SYNTHETIC DATA  \u00b7  EDUCATIONAL DEMONSTRATION',
                   'size': 12, 'color': RGBColor(0xBF, 0xDA, 0xDB),
                   'bold': True, 'font': HEAD}])

# ===========================================================================
# SLIDE 2 - Disclosures
# ===========================================================================
s = add_slide()
top = header(s, 'Disclosures')
_, tf = textbox(s, MARGIN, Inches(3.0), CONTENT_W, Inches(2.0),
                anchor=MSO_ANCHOR.MIDDLE)
para(tf, 'No financial disclosures', 30, INK, bold=True, font=HEAD,
     first=True, align=PP_ALIGN.CENTER, space_after=18)
para(tf, 'No conflicts of interest', 30, INK, bold=True, font=HEAD,
     align=PP_ALIGN.CENTER, space_after=0)
footer(s, 2)

# ===========================================================================
# SLIDE 3 - AI Landscape
# ===========================================================================
s = add_slide()
top = header(s, 'Many models available \u2014 why Claude', kicker='The AI landscape')
# left: models list
lx = MARGIN
_, tf = textbox(s, lx, top, Inches(5.6), Inches(4.6))
para(tf, 'Today\u2019s frontier models', 17, GRAY, bold=True, font=HEAD,
     first=True, space_after=14)
for name, who in [('Claude', 'Anthropic'), ('GPT-4 / o-series', 'OpenAI'),
                  ('Gemini', 'Google'), ('Llama', 'Meta'),
                  ('Mistral, and others', '')]:
    p = tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = '\u2022  ' + name
    _set_font(r, 19, INK, bold=(name == 'Claude'), font=BODY)
    if who:
        r2 = p.add_run(); r2.text = '   ' + who
        _set_font(r2, 14, MUTE, italic=True, font=BODY)
# right: why claude card
rx = Inches(7.1)
cd = card(s, rx, top, Inches(5.5), Inches(4.5), fill=LIGHT, line=RULE)
_, tf = textbox(s, rx + Inches(0.4), top + Inches(0.35), Inches(4.7), Inches(3.9))
para(tf, 'WHY WE FOCUS ON CLAUDE', 14, TEAL, bold=True, font=HEAD,
     first=True, space_after=14)
for t in ['Strong reasoning for research workflows',
          'Built for agentic, multi-step tool use',
          'Large context \u2014 whole datasets & papers',
          'Reliable code generation',
          'It is simply what we use most']:
    p = tf.add_paragraph(); p.space_after = Pt(11)
    r = p.add_run(); r.text = '\u2014  ' + t
    _set_font(r, 16, INK, font=BODY)
footer(s, 3)

# ===========================================================================
# SLIDE 4 - LLM -> Agentic (native two-panel)
# ===========================================================================
s = add_slide()
top = header(s, 'From chatbots to autonomous agents',
             kicker='The shift that changes everything')
panel_y = top + Inches(0.1)
panel_h = Inches(3.7)
pw = Inches(5.5)
gap = SW - 2 * MARGIN - 2 * pw
# Left panel: traditional LLM
lp = card(s, MARGIN, panel_y, pw, panel_h, fill=LIGHT, line=RULE)
_, tf = textbox(s, MARGIN + Inches(0.4), panel_y + Inches(0.3),
                pw - Inches(0.8), panel_h - Inches(0.6))
para(tf, 'TRADITIONAL LLM', 14, GRAY, bold=True, font=HEAD, first=True,
     space_after=4)
para(tf, '~2023', 12, MUTE, italic=True, font=BODY, space_after=14)
for t in ['One question \u2192 one answer',
          'No memory between tasks',
          'You copy, paste, and run the code',
          'You do the work']:
    p = tf.add_paragraph(); p.space_after = Pt(10)
    r = p.add_run(); r.text = '\u2022  ' + t
    _set_font(r, 16, INK, font=BODY)
# Right panel: agentic
rp = card(s, MARGIN + pw + gap, panel_y, pw, panel_h, fill=NAVY, line=NAVY)
_, tf = textbox(s, MARGIN + pw + gap + Inches(0.4), panel_y + Inches(0.3),
                pw - Inches(0.8), panel_h - Inches(0.6))
para(tf, 'AGENTIC SYSTEM', 14, RGBColor(0x8F, 0xC9, 0xCA), bold=True, font=HEAD,
     first=True, space_after=4)
para(tf, '2024\u20132025', 12, RGBColor(0xAD, 0xBE, 0xD0), italic=True, space_after=14)
for t in ['Plans and executes multi-step tasks',
          'Uses tools: code, files, web, terminal',
          'Remembers context and iterates',
          'The agent does the work']:
    p = tf.add_paragraph(); p.space_after = Pt(10)
    r = p.add_run(); r.text = '\u2022  ' + t
    _set_font(r, 16, WHITE, font=BODY)
# center arrow
ar = rect(s, MARGIN + pw + Inches(0.02), panel_y + Inches(1.5),
          gap - Inches(0.04), Inches(0.7), fill=TEAL,
          shape=MSO_SHAPE.RIGHT_ARROW)
# takeaway line
_, tf = textbox(s, MARGIN, panel_y + panel_h + Inches(0.25), CONTENT_W, Inches(0.6))
para(tf, 'Like a research assistant who can code, analyze data, and build tools '
         '\u2014 not just answer questions.', 16, GRAY, italic=True, font=BODY,
     first=True, align=PP_ALIGN.CENTER, space_after=0)
footer(s, 4)

# ===========================================================================
# SLIDE 5 - 4-C Framework (native 2x2 cards)
# ===========================================================================
s = add_slide()
top = header(s, 'The 4-C framework for prompting agents',
             kicker='How to talk to an agent')
cards = [
    ('C', 'Context', 'Set the stage', '"I study CNV after choroidal rupture."', TEAL),
    ('C', 'Clear goal', 'State the outcome', '"Build a model for time-to-CNV."', NAVY),
    ('C', 'Constraints', 'Give specifications', '"1000 patients, 25+ variables, lit-based."', AMBER),
    ('C', 'Collaboration', 'Iterate together', '"Show me the structure first."', CRIMSON),
]
cw, ch = Inches(5.75), Inches(1.95)
gx, gy = Inches(0.35), Inches(0.3)
x0 = MARGIN
for i, (letter, title, sub, ex, col) in enumerate(cards):
    cx = x0 + (i % 2) * (cw + gx)
    cy = top + (i // 2) * (ch + gy)
    cd = card(s, cx, cy, cw, ch, fill=WHITE, line=RULE, line_w=Pt(1.25))
    # color tab
    rect(s, cx, cy, Inches(0.14), ch, fill=col,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = textbox(s, cx + Inches(0.4), cy + Inches(0.22),
                    cw - Inches(0.7), ch - Inches(0.4))
    para(tf, title, 19, NAVY, bold=True, font=HEAD, first=True, space_after=2)
    para(tf, sub, 13, col, bold=True, font=HEAD, space_after=8)
    para(tf, ex, 14, GRAY, italic=True, font=BODY, space_after=0)
footer(s, 5)

# ===========================================================================
# SLIDE 6 - Today's demo
# ===========================================================================
s = add_slide()
top = header(s, 'Predicting CNV after choroidal rupture',
             kicker="Today's demonstration")
# research question banner
qb = card(s, MARGIN, top, CONTENT_W, Inches(1.0), fill=LIGHT, line=RULE)
_, tf = textbox(s, MARGIN + Inches(0.4), top + Inches(0.16),
                CONTENT_W - Inches(0.8), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
para(tf, 'Research question', 13, TEAL, bold=True, font=HEAD, first=True,
     space_after=2)
para(tf, 'Which patients develop CNV \u2014 and how long until it appears?',
     20, INK, bold=True, font=HEAD, space_after=0)
# steps
list_y = top + Inches(1.35)
_, tf = textbox(s, MARGIN + Inches(0.2), list_y, CONTENT_W, Inches(3.0))
para(tf, 'We will build the whole pipeline live, with Claude Code as the agent:',
     15, GRAY, font=BODY, first=True, space_after=12)
for i, t in enumerate([
        'Generate a synthetic database (1000 patients, 33 variables)',
        'Explore the data and run descriptive statistics',
        'Build predictive models (Cox + Random Forest)',
        'Create publication-quality figures',
        'Ship an interactive risk-prediction web app'], 1):
    p = tf.add_paragraph(); p.space_after = Pt(10)
    r = p.add_run(); r.text = f'{i}.  '
    _set_font(r, 17, TEAL, bold=True, font=HEAD)
    r2 = p.add_run(); r2.text = t
    _set_font(r2, 17, INK, font=BODY)
footer(s, 6)

# ===========================================================================
# SLIDE 7 - LIVE DEMO divider
# ===========================================================================
s = add_slide()
gradient_bg(s, CRIMSON, RGBColor(0x8E, 0x24, 0x1A))
_, tf = textbox(s, Inches(1), Inches(2.9), Inches(11.3), Inches(1.0),
                anchor=MSO_ANCHOR.MIDDLE)
para(tf, 'THE DEMONSTRATION', 52, WHITE, bold=True, font=HEAD, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
_, tf = textbox(s, Inches(1), Inches(4.05), Inches(11.3), Inches(0.6))
para(tf, 'A walkthrough of how Claude Code built the pipeline, end to end', 20,
     RGBColor(0xF6, 0xD5, 0xD0), font=BODY, first=True, align=PP_ALIGN.CENTER,
     space_after=0)

# ===========================================================================
# SLIDE 8 - 5-phase pipeline (native steps)
# ===========================================================================
s = add_slide()
top = header(s, 'A five-phase pipeline (~18 minutes)',
             kicker='What the live demo covers')
phases = [
    ('01', 'Data\ngeneration', '5 min', TEAL),
    ('02', 'Exploratory\nanalysis', '4 min', NAVY),
    ('03', 'Predictive\nmodeling', '4 min', AMBER),
    ('04', 'Visualization', '2 min', GREEN),
    ('05', 'Web app\ndeployment', '3 min', CRIMSON),
]
n = len(phases)
gap = Inches(0.3)
total_gap = gap * (n - 1)
pw = (CONTENT_W - total_gap) / n
py = top + Inches(0.5)
ph = Inches(2.7)
for i, (num, name, dur, col) in enumerate(phases):
    px = MARGIN + i * (pw + gap)
    cd = card(s, px, py, pw, ph, fill=WHITE, line=col, line_w=Pt(1.5))
    # number circle
    circ = rect(s, px + pw/2 - Inches(0.45), py + Inches(0.3),
                Inches(0.9), Inches(0.9), fill=col, shape=MSO_SHAPE.OVAL)
    shape_text(circ, [{'text': num, 'size': 22, 'color': WHITE, 'bold': True,
                       'font': HEAD}])
    _, tf = textbox(s, px + Inches(0.1), py + Inches(1.35),
                    pw - Inches(0.2), Inches(1.2), anchor=MSO_ANCHOR.TOP,
                    align=PP_ALIGN.CENTER)
    para(tf, name, 15, NAVY, bold=True, font=HEAD, first=True,
         align=PP_ALIGN.CENTER, space_after=6, line=1.0)
    para(tf, dur, 13, col, bold=True, font=HEAD, align=PP_ALIGN.CENTER,
         space_after=0)
    if i < n - 1:
        ar = rect(s, px + pw + Inches(0.02), py + ph/2 - Inches(0.12),
                  gap - Inches(0.04), Inches(0.24), fill=RULE,
                  shape=MSO_SHAPE.RIGHT_ARROW)
_, tf = textbox(s, MARGIN, py + ph + Inches(0.3), CONTENT_W, Inches(0.5))
para(tf, 'One natural-language prompt kicks off each phase; the agent writes '
         'and runs the code.', 15, GRAY, italic=True, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
footer(s, 8)

# ===========================================================================
# SLIDE 9 - Phase 1: agent-written code (native code panel)
# ===========================================================================
s = add_slide()
top = header(s, 'Phase 1 \u2014 the agent writes the code',
             kicker='You prompt, it builds')
prompt_bar(s, top, 'Generate a synthetic dataset of 1,000 choroidal-rupture '
           'patients with 25+ clinically-grounded variables, and save it as a CSV.',
           accent=TEAL)
top = top + Inches(0.72)
panel = card(s, MARGIN, top, CONTENT_W, Inches(4.25), fill=PANEL, line=PANEL,
             radius=0.04)
# title bar of "editor"
bar = rect(s, MARGIN, top, CONTENT_W, Inches(0.5), fill=RGBColor(0x14, 0x1E, 0x29),
           shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None)
try:
    bar.adjustments[0] = 0.10
except Exception:
    pass
_, tf = textbox(s, MARGIN + Inches(0.35), top + Inches(0.1), Inches(6),
                Inches(0.3))
para(tf, 'generate_synthetic_data.py', 13, RGBColor(0x9F, 0xB3, 0xC4),
     bold=True, font=MONO, first=True, space_after=0)
# code lines
code = [
    ('# Agent-written  \u00b7  generate 1000 synthetic patients', 'cmt'),
    ('import pandas as pd, numpy as np', 'kw'),
    ('N = 1000', 'p'),
    ('data = {', 'p'),
    ('    "patient_id": [f"CR{i:04d}" for i in range(1, N+1)],', 'p'),
    ('    "age": np.clip(np.random.normal(35, 15, N), 10, 85),', 'p'),
    ('    "rupture_length_mm": np.random.gamma(2, 1.5, N),', 'p'),
    ('    "foveal_involvement": np.random.choice(["Yes","No"], N),', 'p'),
    ('    # ... 25+ more clinically-grounded variables', 'cmt'),
    ('}', 'p'),
    ('# risk rises with foveal involvement & rupture size', 'cmt'),
    ('cnv_risk = 0.25 + 0.25*foveal + 0.03*rupture_mm', 'p'),
    ('pd.DataFrame(data).to_csv("choroidal_rupture_data.csv")', 'p'),
]
_, tf = textbox(s, MARGIN + Inches(0.35), top + Inches(0.72),
                CONTENT_W - Inches(0.7), Inches(3.6))
for i, (line, kind) in enumerate(code):
    col = (RGBColor(0x6F, 0x9B, 0x76) if kind == 'cmt'
           else RGBColor(0x82, 0xC8, 0xE6) if kind == 'kw'
           else RGBColor(0xE6, 0xED, 0xF3))
    para(tf, line if line else '\u00a0', 12.5, col, font=MONO,
         first=(i == 0), space_after=2, line=1.05)
footer(s, 9)

# ===========================================================================
# SLIDE 10 - Data output (native table)
# ===========================================================================
s = add_slide()
top = header(s, '1000 patients \u00d7 33 variables \u2014 a wide dataset',
             kicker='The generated database', accent=GREEN)
# Framing line: emphasise the WIDTH (every patient carries 33 columns).
_, tf = textbox(s, MARGIN, top, CONTENT_W, Inches(0.4))
para(tf, 'Every patient row carries 33 columns \u2014 grouped into six clinical '
     'domains, from demographics to OCT biomarkers to long-term outcomes.',
     13, GRAY, font=BODY, first=True, space_after=0)

SLATE = RGBColor(0x44, 0x55, 0x66)
groups = [
    ('Demographics & history', NAVY,
     ['Patient ID', 'Age', 'Sex', 'Race', 'Occupational risk',
      'Injury mechanism']),
    ('Rupture & exam', TEAL,
     ['Rupture clock-hour', 'Rupture length', 'Distance from fovea',
      'Foveal involvement', 'Choroidal thickness', 'Time to presentation',
      'Subretinal hemorrhage']),
    ('OCT biomarkers', AMBER,
     ['Baseline CST', 'Outer-retinal disruption', 'RPE disruption',
      'Subfoveal choroid', 'Subretinal fluid', 'Intraretinal cysts',
      'EZ integrity', 'ELM intact', 'OCT quality score']),
    ('Visual acuity', GREEN,
     ['Baseline VA (logMAR)', 'Final VA (logMAR)']),
    ('Treatment', CRIMSON,
     ['Treatment approach', 'Anti-VEGF agent', 'Number of injections',
      'Supplemental laser']),
    ('Outcomes', SLATE,
     ['CNV developed', 'Time to CNV', 'CNV type', 'Treatment response',
      'Follow-up duration']),
]

# 3-column x 2-row grid of category cards listing all 33 variable names.
panel_top = top + Inches(0.55)
ncol, nrow = 3, 2
hgap, vgap = Inches(0.3), Inches(0.25)
card_w = (CONTENT_W - hgap * (ncol - 1)) / ncol
card_h = Inches(2.0)
for idx, (gname, gcol, varnames) in enumerate(groups):
    r, c = divmod(idx, ncol)
    cx = MARGIN + c * (card_w + hgap)
    cy = panel_top + r * (card_h + vgap)
    card(s, cx, cy, card_w, card_h, fill=LIGHT, line=RULE)
    # colored accent rule under the category title
    _, tf = textbox(s, cx + Inches(0.18), cy + Inches(0.13),
                    card_w - Inches(0.36), card_h - Inches(0.24))
    para(tf, f'{gname}  ({len(varnames)})', 12.5, gcol, bold=True,
         font=HEAD, first=True, space_after=5)
    for v in varnames:
        para(tf, v, 9.5, INK, font=BODY, space_after=1, line=1.0)
footer(s, 10)


# ---- helper: full-bleed figure slide -------------------------------------
# (defined here so the EDA slide below can use it)
def figure_slide(num, kicker, title, img, takeaway, accent=TEAL,
                 img_h=Inches(4.55), prompt=None):
    s = add_slide()
    top = header(s, title, kicker=kicker, accent=accent)
    img_top = top + Inches(0.05)
    eff_h = img_h
    if prompt:
        prompt_bar(s, top, prompt, accent=accent)
        img_top = top + Inches(0.72)
        eff_h = min(img_h, Inches(3.8))
    if os.path.exists(f'{BASE}/{img}'):
        from PIL import Image
        iw, ih = Image.open(f'{BASE}/{img}').size
        ratio = iw / ih
        h = eff_h
        w = Emu(int(h * ratio))
        if w > CONTENT_W:
            w = CONTENT_W
            h = Emu(int(w / ratio))
        x = MARGIN + (CONTENT_W - w) // 2
        s.shapes.add_picture(f'{BASE}/{img}', x, img_top, width=w, height=h)
    if takeaway:
        _, tf = textbox(s, MARGIN, Inches(6.5), CONTENT_W, Inches(0.5))
        para(tf, takeaway, 14, GRAY, italic=True, first=True,
             align=PP_ALIGN.CENTER, space_after=0)
    footer(s, num)
    return s


# SLIDE 11 - Exploratory data analysis (raw-data trends, before modeling)
figure_slide(11, 'Phase 2 \u00b7 exploratory analysis',
             'First, explore the data \u2014 the trends jump out',
             'figure_eda_clean.png',
             'Before any model: foveal involvement, larger ruptures and proximity '
             'to the fovea all track with CNV \u2014 most cases appear within ~6 months.',
             accent=NAVY, img_h=Inches(4.5),
             prompt='Before any modeling, explore the data \u2014 which factors track '
                    'with CNV? Give me one clean 4-panel figure.')
# SLIDE 12 - KM
figure_slide(12, 'Phase 3 \u00b7 survival analysis',
             'Foveal involvement drives time to CNV', 'figure1_clean.png',
             'Patients with foveal involvement reach 50% CNV by ~6 months; the '
             'separation is large and highly significant (log-rank p < 0.001).',
             prompt='Run a Kaplan\u2013Meier survival analysis comparing eyes with vs. '
                    'without foveal involvement, and plot time-to-CNV.')
# SLIDE 13 - feature importance
figure_slide(13, 'Phase 3 \u00b7 model drivers',
             'What predicts CNV development', 'figure2_clean.png',
             'Distance from fovea, baseline acuity, and foveal involvement '
             'dominate the random-forest model.',
             prompt='Fit a random forest to predict CNV and show me which variables '
                    'matter most.')
# SLIDE 14 - ROC
figure_slide(14, 'Phase 3 \u00b7 discrimination',
             'Model performance', 'figure3_clean.png',
             'On held-out data, the model ranks a true-CNV eye above a CNV-free '
             'eye ~77% of the time (AUC 0.77; chance = 0.50).',
             img_h=Inches(4.5),
             prompt='How well does the model tell future-CNV eyes from CNV-free '
                    'eyes? Plot the ROC curve and report the AUC.')
# SLIDE 15 - risk stratification
figure_slide(15, 'Phase 4 \u00b7 clinical translation',
             'Risk groups behave as expected', 'figure4_clean.png',
             'Predicted risk tiers map cleanly onto observed CNV rates \u2014 the '
             'basis for a triage tool.', img_h=Inches(4.5),
             prompt='Sort patients into low / medium / high risk tiers and check the '
                    'observed CNV rate in each \u2014 does it line up?')

# ===========================================================================
# SLIDE 15 - Interactive clinical tool
# ===========================================================================
s = add_slide()
top = header(s, 'From model to bedside tool', kicker='Phase 5 \u00b7 the web app',
             accent=GREEN)
prompt_bar(s, top, 'Now turn this into an interactive web app: a clinician enters a '
           'patient and sees their CNV risk, survival curve, and top drivers.',
           accent=GREEN)
top = top + Inches(0.72)
# features column
_, tf = textbox(s, MARGIN, top, Inches(5.4), Inches(3.2))
para(tf, 'A deployable app, built in the same session', 17, NAVY, bold=True,
     font=HEAD, first=True, space_after=11)
for t in ['Live risk gauge \u2014 updates as you type',
          'Individualized vs population survival curve',
          'Honest per-patient drivers (vs a typical case)',
          'Cohort context + downloadable report',
          'Open on any phone on the same wifi']:
    p = tf.add_paragraph(); p.space_after = Pt(8)
    r = p.add_run(); r.text = '\u2014  ' + t
    _set_font(r, 15.5, INK, font=BODY)
# launch strip: QR + URL
qr_y = top + Inches(3.05)
import os as _os
if _os.path.exists('app_qr.png'):
    s.shapes.add_picture('app_qr.png', MARGIN, qr_y, height=Inches(1.25))
_, tf = textbox(s, MARGIN + Inches(1.45), qr_y + Inches(0.1), Inches(4.4), Inches(1.2))
para(tf, 'SCAN TO OPEN', 12, GREEN, bold=True, font=HEAD, first=True, space_after=2)
para(tf, 'run_app.sh  \u2192  http://<wifi-ip>:8501', 12.5, NAVY, bold=True,
     font=MONO, space_after=8)
para(tf, 'CODE ON GITHUB', 10.5, GREEN, bold=True, font=HEAD, space_after=2)
para(tf, 'github.com/Mak-Djulbegovic/cnv-risk-prediction-demo', 10, NAVY,
     font=MONO, space_after=0,
     link='https://github.com/Mak-Djulbegovic/cnv-risk-prediction-demo')
# two case cards
case_x = Inches(6.6)
cw = Inches(6.0)
# low risk
lc = card(s, case_x, top, cw, Inches(1.8), fill=LIGHT, line=GREEN, line_w=Pt(1.5))
_, tf = textbox(s, case_x + Inches(0.4), top + Inches(0.2), cw - Inches(0.8), Inches(1.5))
para(tf, 'CASE 1 \u00b7 LOW RISK', 13, GREEN, bold=True, font=HEAD, first=True, space_after=4)
para(tf, '65 y \u00b7 peripheral rupture \u00b7 fovea spared', 14, GRAY, font=BODY, space_after=8)
para(tf, '25% CNV probability \u00b7 >36-mo median', 20, GREEN, bold=True, font=HEAD, space_after=0)
# high risk
hc = card(s, case_x, top + Inches(2.05), cw, Inches(1.8), fill=LIGHT,
          line=CRIMSON, line_w=Pt(1.5))
_, tf = textbox(s, case_x + Inches(0.4), top + Inches(2.25), cw - Inches(0.8), Inches(1.5))
para(tf, 'CASE 2 \u00b7 HIGH RISK', 13, CRIMSON, bold=True, font=HEAD, first=True, space_after=4)
para(tf, '28 y \u00b7 large rupture \u00b7 foveal + hemorrhage', 14, GRAY, font=BODY, space_after=8)
para(tf, '84% CNV probability \u00b7 6-mo median', 20, CRIMSON, bold=True, font=HEAD, space_after=0)
footer(s, 15)

# ===========================================================================
# SLIDE 16 - Key findings
# ===========================================================================
s = add_slide()
top = header(s, 'What the demo produced', kicker='Key findings')
colw = Inches(5.75)
# left card - dataset / risk factors
lc = card(s, MARGIN, top, colw, Inches(4.4), fill=LIGHT, line=RULE)
_, tf = textbox(s, MARGIN + Inches(0.4), top + Inches(0.3), colw - Inches(0.8), Inches(3.9))
para(tf, 'DATASET & RISK FACTORS', 14, TEAL, bold=True, font=HEAD, first=True, space_after=12)
for t in ['1000 patients \u00b7 33 variables',
          '634 CNV cases (63.4%)',
          'Median time to CNV: 6.7 months']:
    p = tf.add_paragraph(); p.space_after = Pt(9)
    r = p.add_run(); r.text = '\u2022  ' + t; _set_font(r, 16, INK, font=BODY)
para(tf, 'Top risk factors', 14, NAVY, bold=True, font=HEAD, space_before=8, space_after=8)
for t in ['Foveal involvement (HR 1.79)',
          'Distance from fovea',
          'Outer retinal disruption',
          'Baseline visual acuity']:
    p = tf.add_paragraph(); p.space_after = Pt(7)
    r = p.add_run(); r.text = '\u2022  ' + t; _set_font(r, 15, INK, font=BODY)
# right card - performance / impact
rc = card(s, MARGIN + colw + Inches(0.35), top, colw, Inches(4.4), fill=LIGHT, line=RULE)
rx = MARGIN + colw + Inches(0.35) + Inches(0.4)
_, tf = textbox(s, rx, top + Inches(0.3), colw - Inches(0.8), Inches(3.9))
para(tf, 'MODEL PERFORMANCE', 14, GREEN, bold=True, font=HEAD, first=True, space_after=12)
for t in ['Cox C-index: 0.66',
          'Random Forest AUC: 0.77',
          'Sensitivity 85% \u00b7 Specificity 53%']:
    p = tf.add_paragraph(); p.space_after = Pt(9)
    r = p.add_run(); r.text = '\u2022  ' + t; _set_font(r, 16, INK, font=BODY)
para(tf, 'Why it matters', 14, NAVY, bold=True, font=HEAD, space_before=8, space_after=8)
for t in ['Risk-stratified monitoring',
          'Earlier intervention for high-risk eyes',
          'Smarter use of clinic resources']:
    p = tf.add_paragraph(); p.space_after = Pt(7)
    r = p.add_run(); r.text = '\u2022  ' + t; _set_font(r, 15, INK, font=BODY)
footer(s, 16)

# ===========================================================================
# SLIDE 17 - Beyond this demo (3 columns)
# ===========================================================================
s = add_slide()
top = header(s, 'Where else this applies', kicker='Beyond the demo')
cols = [
    ('Research', TEAL, ['Systematic reviews', 'Meta-analyses', 'Trial design',
                        'Power calculations', 'Grant writing']),
    ('Clinical', NAVY, ['EHR data extraction', 'Clinical calculators',
                        'Prediction tools', 'QI dashboards', 'Outcome tracking']),
    ('Education', AMBER, ['Teaching cases', 'Learning modules',
                          'Competency tests', 'Study materials', 'Simulations']),
]
cw = Inches(3.9)
cgap = (CONTENT_W - cw * 3) / 2
for i, (title, col, items) in enumerate(cols):
    cx = MARGIN + i * (cw + cgap)
    cd = card(s, cx, top, cw, Inches(4.2), fill=WHITE, line=RULE, line_w=Pt(1.25))
    rect(s, cx, top, cw, Inches(0.7), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, top + Inches(0.35), cw, Inches(0.35), fill=col)  # square bottom of header
    _, tf = textbox(s, cx, top + Inches(0.12), cw, Inches(0.5), align=PP_ALIGN.CENTER)
    para(tf, title, 18, WHITE, bold=True, font=HEAD, first=True, align=PP_ALIGN.CENTER, space_after=0)
    _, tf = textbox(s, cx + Inches(0.4), top + Inches(1.0), cw - Inches(0.7), Inches(3.0))
    first = True
    for it in items:
        para(tf, '\u2014  ' + it, 15, INK, font=BODY, first=first, space_after=11)
        first = False
footer(s, 17)

# ===========================================================================
# SLIDE 18 - Practical guidance (do / don't)
# ===========================================================================
s = add_slide()
top = header(s, 'Using agents well', kicker='Practical guidance')
cw = Inches(5.75)
# DO card
dc = card(s, MARGIN, top, cw, Inches(4.4), fill=LIGHT, line=GREEN, line_w=Pt(1.5))
_, tf = textbox(s, MARGIN + Inches(0.4), top + Inches(0.3), cw - Inches(0.8), Inches(3.9))
para(tf, 'DO', 16, GREEN, bold=True, font=HEAD, first=True, space_after=12)
for t in ['Use the 4-C framework for clear prompts',
          'Let the agent work autonomously',
          'Iterate with feedback on its output',
          'Always verify the logic and numbers',
          'Use it for exploration & prototyping']:
    p = tf.add_paragraph(); p.space_after = Pt(11)
    r = p.add_run(); r.text = '\u2713  ' + t; _set_font(r, 16, INK, font=BODY)
# DON'T card
xc = card(s, MARGIN + cw + Inches(0.35), top, cw, Inches(4.4), fill=LIGHT,
          line=CRIMSON, line_w=Pt(1.5))
xx = MARGIN + cw + Inches(0.35) + Inches(0.4)
_, tf = textbox(s, xx, top + Inches(0.3), cw - Inches(0.8), Inches(3.9))
para(tf, 'DON\u2019T', 16, CRIMSON, bold=True, font=HEAD, first=True, space_after=12)
for t in ['Skip biostatistician review of final work',
          'Treat output as publication-ready',
          'Replace clinical judgment',
          'Feed real PHI without IRB & de-identification',
          'Trust numbers you have not checked']:
    p = tf.add_paragraph(); p.space_after = Pt(11)
    r = p.add_run(); r.text = '\u2717  ' + t; _set_font(r, 16, INK, font=BODY)
footer(s, 18)

# ===========================================================================
# SLIDE 19 - Resources & considerations
# ===========================================================================
s = add_slide()
top = header(s, 'Getting started & guardrails', kicker='Resources')
cw = Inches(5.75)
lc = card(s, MARGIN, top, cw, Inches(4.4), fill=LIGHT, line=RULE)
_, tf = textbox(s, MARGIN + Inches(0.4), top + Inches(0.3), cw - Inches(0.8), Inches(3.9))
para(tf, 'TRY IT', 14, TEAL, bold=True, font=HEAD, first=True, space_after=12)
for t in ['Claude.ai \u2014 chat, free tier available',
          'Claude Code \u2014 agentic CLI workflows',
          'Cursor / Windsurf \u2014 AI-native editors',
          'GitHub Copilot \u2014 in-editor assist']:
    p = tf.add_paragraph(); p.space_after = Pt(12)
    r = p.add_run(); r.text = '\u2014  ' + t; _set_font(r, 16, INK, font=BODY)
rc = card(s, MARGIN + cw + Inches(0.35), top, cw, Inches(4.4), fill=LIGHT,
          line=AMBER, line_w=Pt(1.5))
rx = MARGIN + cw + Inches(0.35) + Inches(0.4)
_, tf = textbox(s, rx, top + Inches(0.3), cw - Inches(0.8), Inches(3.9))
para(tf, 'BEFORE YOU USE REAL DATA', 14, AMBER, bold=True, font=HEAD, first=True, space_after=12)
for t in ['Institutional AI policies', 'HIPAA & data privacy',
          'IRB approval & de-identification', 'Authorship guidelines',
          'Verify every output']:
    p = tf.add_paragraph(); p.space_after = Pt(11)
    r = p.add_run(); r.text = '\u2022  ' + t; _set_font(r, 16, INK, font=BODY)
footer(s, 19)

# ===========================================================================
# SLIDE 20 - Thank you / Q&A
# ===========================================================================
s = add_slide()
gradient_bg(s, RGBColor(0x0C, 0x1E, 0x33), NAVY)
rect(s, 0, Inches(3.4), SW, Pt(2.5), fill=TEAL)
_, tf = textbox(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.0))
para(tf, 'Questions & Discussion', 44, WHITE, bold=True, font=HEAD, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
_, tf = textbox(s, Inches(1), Inches(3.65), Inches(11.3), Inches(0.6))
para(tf, 'From research question to clinical tool \u2014 in one session', 19,
     RGBColor(0xC7, 0xD3, 0xE0), font=BODY, first=True, align=PP_ALIGN.CENTER,
     space_after=0)
chip = card(s, Inches(3.67), Inches(5.5), Inches(6.0), Inches(0.95),
            fill=RGBColor(0x1C, 0x33, 0x4D), line=TEAL, line_w=Pt(1))
_, tf = textbox(s, Inches(3.67), Inches(5.62), Inches(6.0), Inches(0.75),
                anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
para(tf, 'This presentation was built with Claude', 16, WHITE, bold=True,
     font=HEAD, first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, 'Synthetic data \u00b7 Educational purposes only', 12,
     RGBColor(0xAD, 0xBE, 0xD0), italic=True, font=BODY, align=PP_ALIGN.CENTER,
     space_after=0)

# ---- save -----------------------------------------------------------------
out = f'{BASE}/AI_Ophthalmology_Presentation.pptx'
prs.save(out)
print(f'Saved {out} with {len(prs.slides._sldIdLst)} slides')
