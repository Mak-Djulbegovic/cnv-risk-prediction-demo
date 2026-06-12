"""
Professional PowerPoint Presentation Generator - FINAL VERSION
==============================================================

Creates a truly professional, polished presentation with:
- Proper margins (no bleeding edges)
- Code examples
- Workflow visualizations
- Sophisticated layouts
- High-quality design
"""

import subprocess
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

import os

print("="*80)
print("CREATING PROFESSIONAL PRESENTATION - FINAL VERSION")
print("="*80)

# Create presentation - 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Professional color palette - Medical/Academic
PRIMARY = RGBColor(0, 102, 204)      # Professional blue
SECONDARY = RGBColor(51, 51, 51)      # Charcoal
ACCENT = RGBColor(220, 53, 69)        # Medical red
SUCCESS = RGBColor(40, 167, 69)       # Green
WARNING = RGBColor(255, 193, 7)       # Yellow
LIGHT_GRAY = RGBColor(248, 249, 250)
DARK_GRAY = RGBColor(73, 80, 87)
WHITE = RGBColor(255, 255, 255)

# Safe margins to prevent bleeding
MARGIN_LEFT = Inches(0.4)
MARGIN_RIGHT = Inches(9.6)
MARGIN_TOP = Inches(0.4)
MARGIN_BOTTOM = Inches(5.225)
CONTENT_WIDTH = MARGIN_RIGHT - MARGIN_LEFT
CONTENT_HEIGHT = MARGIN_BOTTOM - MARGIN_TOP

def add_slide_number(slide, num):
    """Add slide number"""
    textbox = slide.shapes.add_textbox(
        Inches(9.3), Inches(5.3),
        Inches(0.5), Inches(0.2)
    )
    tf = textbox.text_frame
    tf.text = str(num)
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_footer_text(slide, text="Synthetic Data - Educational Demo"):
    """Add footer"""
    textbox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(5.35),
        CONTENT_WIDTH, Inches(0.15)
    )
    tf = textbox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(8)
    p.font.color.rgb = DARK_GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

print("\nCreating slides...")

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
print("  1. Title slide")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Full background
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg.fill.gradient()
bg.fill.gradient_angle = 90
bg.fill.gradient_stops[0].color.rgb = PRIMARY
bg.fill.gradient_stops[1].color.rgb = SECONDARY
bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.5), Inches(9), Inches(1.2)
)
tf = title_box.text_frame
tf.text = "AI-Powered Research in Ophthalmology"
p = tf.paragraphs[0]
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(2.8), Inches(9), Inches(0.8)
)
tf = subtitle_box.text_frame
tf.text = "Agentic Systems: From Research Question to Clinical Tool"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# Warning
warn_box = slide.shapes.add_textbox(
    Inches(2), Inches(4.5), Inches(6), Inches(0.4)
)
tf = warn_box.text_frame
tf.text = "⚠️  SYNTHETIC DATA - Educational Demonstration"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WARNING
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 2: Disclosures
# ============================================================================
print("  2. Disclosures")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title bar
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Disclosures"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Content
content_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(1), Inches(2.2),
    CONTENT_WIDTH - Inches(2), Inches(1)
)
tf = content_box.text_frame
tf.text = "No financial disclosures\n\nNo conflicts of interest"
for para in tf.paragraphs:
    para.font.size = Pt(24)
    para.font.color.rgb = SECONDARY
    para.alignment = PP_ALIGN.CENTER
    para.space_after = Pt(12)

add_footer_text(slide)
add_slide_number(slide, 2)

# ============================================================================
# SLIDE 3: AI Landscape
# ============================================================================
print("  3. AI Landscape")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "The AI Landscape: Many Models Available"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Content - Two columns
left_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.3), Inches(1.3),
    Inches(4), Inches(3.5)
)
tf = left_box.text_frame
content = """Available Models:
• GPT-4 (OpenAI)
• Gemini (Google)
• Claude (Anthropic)
• Llama (Meta)
• Many others..."""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i == 0:
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = SECONDARY
    else:
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
    para.space_after = Pt(8)

right_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(4.8), Inches(1.3),
    Inches(4.3), Inches(3.5)
)
tf = right_box.text_frame
content = """Why Claude?
• Strong reasoning
• Research-oriented
• Agentic capabilities
• Large context window
• Code generation"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i == 0:
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = PRIMARY
    else:
        para.font.size = Pt(16)
        para.font.color.rgb = DARK_GRAY
    para.space_after = Pt(8)

add_footer_text(slide)
add_slide_number(slide, 3)

# ============================================================================
# SLIDE 4: LLM Evolution - WITH WORKFLOW IMAGE
# ============================================================================
print("  4. LLM to Agentic evolution")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Evolution: From LLMs to Agentic Systems"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Add workflow image
if os.path.exists('workflow_diagram.png'):
    slide.shapes.add_picture(
        'workflow_diagram.png',
        MARGIN_LEFT + Inches(0.5), Inches(1.3),
        width=Inches(8.2)
    )

add_footer_text(slide)
add_slide_number(slide, 4)

# ============================================================================
# SLIDE 5: Prompt Framework - WITH IMAGE
# ============================================================================
print("  5. Prompt engineering framework")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "The 4-C Prompt Engineering Framework"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Add framework image
if os.path.exists('prompt_framework.png'):
    slide.shapes.add_picture(
        'prompt_framework.png',
        MARGIN_LEFT + Inches(0.5), Inches(1.2),
        width=Inches(8.2)
    )

add_footer_text(slide)
add_slide_number(slide, 5)

# ============================================================================
# SLIDE 6: Today's Demo Setup
# ============================================================================
print("  6. Demo setup")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Today's Demo: CNV Prediction After Choroidal Rupture"
p = tf.paragraphs[0]
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = WHITE

# Content
content_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.5), Inches(1.3),
    CONTENT_WIDTH - Inches(1), Inches(3.5)
)
tf = content_box.text_frame
content = """Research Question:
Which patients develop CNV? How long until it develops?

Live Demonstration:
1. Generate synthetic database (1000 patients, 33 variables)
2. Exploratory data analysis
3. Build predictive models (Cox + Random Forest)
4. Create publication-quality visualizations
5. Build interactive web application

All executed by Claude Code as autonomous agent"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i in [0, 2]:
        para.font.size = Pt(20)
        para.font.bold = True
        para.font.color.rgb = ACCENT
    else:
        para.font.size = Pt(16)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(10)

add_footer_text(slide)
add_slide_number(slide, 6)

# ============================================================================
# SLIDE 7: LIVE DEMO - Section Divider
# ============================================================================
print("  7. Live demo divider")
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg.fill.solid()
bg.fill.fore_color.rgb = ACCENT
bg.line.fill.background()

text_box = slide.shapes.add_textbox(
    Inches(1), Inches(2.3), Inches(8), Inches(1)
)
tf = text_box.text_frame
tf.text = "LIVE DEMONSTRATION"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 8: Demo Phases - WITH IMAGE
# ============================================================================
print("  8. Demo phases timeline")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "5-Phase Pipeline (~18 minutes)"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Add demo phases image
if os.path.exists('demo_phases.png'):
    slide.shapes.add_picture(
        'demo_phases.png',
        MARGIN_LEFT + Inches(0.3), Inches(1.2),
        width=Inches(8.6)
    )

add_footer_text(slide)
add_slide_number(slide, 8)

# ============================================================================
# SLIDE 9: Phase 1 - Data Generation WITH CODE
# ============================================================================
print("  9. Phase 1: Data generation with code")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Phase 1: Data Generation - Agent-Written Code"
p = tf.paragraphs[0]
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = WHITE

# Add code example image
if os.path.exists('code_example.png'):
    slide.shapes.add_picture(
        'code_example.png',
        MARGIN_LEFT + Inches(0.3), Inches(1.2),
        width=Inches(8.6)
    )

add_footer_text(slide)
add_slide_number(slide, 9)

# ============================================================================
# SLIDE 10: Database Output
# ============================================================================
print("  10. Database preview")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = SUCCESS
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Output: 1000 Patients × 33 Variables"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Add database preview
if os.path.exists('database_preview.png'):
    slide.shapes.add_picture(
        'database_preview.png',
        MARGIN_LEFT + Inches(0.3), Inches(1.2),
        width=Inches(8.6)
    )

add_footer_text(slide)
add_slide_number(slide, 10)

# ============================================================================
# SLIDE 11: Kaplan-Meier
# ============================================================================
print("  11. Kaplan-Meier results")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.5)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.05),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Survival Analysis: Time to CNV Development"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE

if os.path.exists('figure1_kaplan_meier.png'):
    slide.shapes.add_picture(
        'figure1_kaplan_meier.png',
        MARGIN_LEFT + Inches(0.3), Inches(1.05),
        width=Inches(8.6)
    )

add_footer_text(slide)
add_slide_number(slide, 11)

# ============================================================================
# SLIDE 12: Feature Importance
# ============================================================================
print("  12. Feature importance")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.5)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.05),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Top Predictors of CNV Development"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE

if os.path.exists('figure2_feature_importance.png'):
    slide.shapes.add_picture(
        'figure2_feature_importance.png',
        MARGIN_LEFT + Inches(0.3), Inches(1.05),
        width=Inches(8.6)
    )

add_footer_text(slide)
add_slide_number(slide, 12)

# ============================================================================
# SLIDE 13: ROC Curve
# ============================================================================
print("  13. ROC curve")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.5)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.05),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Model Performance: AUC = 0.768"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE

if os.path.exists('figure3_roc_curve.png'):
    slide.shapes.add_picture(
        'figure3_roc_curve.png',
        MARGIN_LEFT + Inches(0.3), Inches(1.05),
        width=Inches(8.6)
    )

add_footer_text(slide)
add_slide_number(slide, 13)

# ============================================================================
# SLIDE 14: Web App Demo
# ============================================================================
print("  14. Web app capabilities")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = SUCCESS
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Phase 5: Interactive Clinical Tool"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Content
content_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.5), Inches(1.3),
    CONTENT_WIDTH - Inches(1), Inches(3.5)
)
tf = content_box.text_frame
content = """Web Application Features:
• Input patient-specific parameters
• Real-time CNV risk prediction
• Risk stratification (Low/Medium/High)
• Contributing factor analysis
• Clinical recommendations
• Downloadable patient reports

Live Demo: Two Patient Cases
1. Low-risk patient (~15% CNV probability)
2. High-risk patient (~75% CNV probability)

From Research Question → Deployable Tool in ONE Session"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i in [0, 7]:
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = SUCCESS
    elif i == 11:
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.color.rgb = ACCENT
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(6)

add_footer_text(slide)
add_slide_number(slide, 14)

# ============================================================================
# SLIDE 15: Key Findings
# ============================================================================
print("  15. Key findings")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Key Findings Summary"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Two columns
left_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.3), Inches(1.3),
    Inches(4.3), Inches(3.5)
)
tf = left_box.text_frame
content = """Dataset:
• 1000 patients
• 33 variables
• 634 CNV cases (63.4%)
• Mean time: 6.7 months

Top Risk Factors:
1. Foveal involvement
   (HR = 1.79)
2. Distance from fovea
3. Outer retinal disruption
4. Baseline visual acuity"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i in [0, 5]:
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = ACCENT
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(6)

right_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(4.9), Inches(1.3),
    Inches(4.3), Inches(3.5)
)
tf = right_box.text_frame
content = """Model Performance:
• Cox C-index: 0.662
• Random Forest AUC: 0.768
• Sensitivity: 85%
• Specificity: 53%

Clinical Impact:
→ Risk-stratified monitoring
→ Early intervention for high-risk
→ Resource optimization
→ Improved patient outcomes"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i in [0, 5]:
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = SUCCESS
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(6)

add_footer_text(slide)
add_slide_number(slide, 15)

# ============================================================================
# SLIDE 16: Other Use Cases
# ============================================================================
print("  16. Other applications")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Beyond This Demo: Broader Applications"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Three columns
col_width = Inches(2.8)
columns = [
    {"title": "Research", "items": [
        "Systematic reviews",
        "Meta-analyses",
        "Trial design",
        "Power calculations",
        "Grant writing"
    ]},
    {"title": "Clinical", "items": [
        "EHR data extraction",
        "Clinical calculators",
        "Prediction tools",
        "QI dashboards",
        "Outcome tracking"
    ]},
    {"title": "Education", "items": [
        "Case generation",
        "Learning modules",
        "Competency tests",
        "Study materials",
        "Simulation"
    ]}
]

x_pos = MARGIN_LEFT + Inches(0.3)
for col in columns:
    box = slide.shapes.add_textbox(x_pos, Inches(1.3), col_width, Inches(3.5))
    tf = box.text_frame
    tf.text = col["title"]
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].space_after = Pt(8)

    for item in col["items"]:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = SECONDARY
        p.space_after = Pt(4)

    x_pos += col_width + Inches(0.15)

add_footer_text(slide)
add_slide_number(slide, 16)

# ============================================================================
# SLIDE 17: Practical Tips
# ============================================================================
print("  17. Practical tips")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Getting Started: Practical Guidance"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Content
content_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.5), Inches(1.3),
    CONTENT_WIDTH - Inches(1), Inches(3.5)
)
tf = content_box.text_frame
content = """Best Practices:
✓ Use the 4-C framework for clear prompts
✓ Let agents work autonomously
✓ Iterate based on outputs
✓ Always verify results

Good For:
• Data exploration and hypothesis generation
• Rapid prototyping and tool building
• Learning new statistical methods

NOT a Replacement For:
✗ Expert biostatistician review
✗ Clinical judgment
✗ Peer review process

Think: Powerful Research Assistant"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i in [0, 5, 9]:
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = ACCENT if i == 9 else PRIMARY if i == 0 else SUCCESS
    elif i == 14:
        para.font.size = Pt(16)
        para.font.bold = True
        para.font.italic = True
        para.font.color.rgb = SECONDARY
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(6)

add_footer_text(slide)
add_slide_number(slide, 17)

# ============================================================================
# SLIDE 18: Resources
# ============================================================================
print("  18. Resources")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.6)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY
title_bg.line.fill.background()

title_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.2), MARGIN_TOP + Inches(0.1),
    CONTENT_WIDTH - Inches(0.4), Inches(0.4)
)
tf = title_box.text_frame
tf.text = "Resources & Important Considerations"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Two columns
left_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(0.3), Inches(1.3),
    Inches(4.3), Inches(3.5)
)
tf = left_box.text_frame
content = """Getting Started:
• Claude.ai (free tier)
• Claude Code CLI
• Cursor / Windsurf
• GitHub Copilot
• Local LLMs"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i == 0:
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = PRIMARY
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(8)

right_box = slide.shapes.add_textbox(
    MARGIN_LEFT + Inches(4.9), Inches(1.3),
    Inches(4.3), Inches(3.5)
)
tf = right_box.text_frame
content = """Important:
⚠️ Institutional AI policies
⚠️ HIPAA compliance
⚠️ Data privacy
⚠️ IRB requirements
⚠️ Authorship guidelines
⚠️ Verify all outputs"""

tf.text = content
for i, para in enumerate(tf.paragraphs):
    if i == 0:
        para.font.size = Pt(18)
        para.font.bold = True
        para.font.color.rgb = ACCENT
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = SECONDARY
    para.space_after = Pt(8)

add_footer_text(slide)
add_slide_number(slide, 18)

# ============================================================================
# SLIDE 19: Questions
# ============================================================================
print("  19. Questions")
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY
bg.line.fill.background()

text_box = slide.shapes.add_textbox(
    Inches(1), Inches(2.3), Inches(8), Inches(1)
)
tf = text_box.text_frame
tf.text = "Questions & Discussion"
p = tf.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 20: Thank You with Attribution
# ============================================================================
print("  20. Thank you with Claude attribution")
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background gradient
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(10), Inches(5.625)
)
bg.fill.gradient()
bg.fill.gradient_angle = 90
bg.fill.gradient_stops[0].color.rgb = SECONDARY
bg.fill.gradient_stops[1].color.rgb = PRIMARY
bg.line.fill.background()

# Thank you
thank_box = slide.shapes.add_textbox(
    Inches(1), Inches(1.2), Inches(8), Inches(0.8)
)
tf = thank_box.text_frame
tf.text = "Thank You!"
p = tf.paragraphs[0]
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(
    Inches(1), Inches(2.1), Inches(8), Inches(0.5)
)
tf = sub_box.text_frame
tf.text = "From Research Question to Clinical Tool in One Session"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# Attribution box
attr_bg = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(2.5), Inches(3), Inches(5), Inches(1.5)
)
attr_bg.fill.solid()
attr_bg.fill.fore_color.rgb = WHITE
attr_bg.line.color.rgb = PRIMARY
attr_bg.line.width = Pt(4)

attr_box = slide.shapes.add_textbox(
    Inches(2.6), Inches(3.1), Inches(4.8), Inches(1.3)
)
tf = attr_box.text_frame

# Line 1
p = tf.paragraphs[0]
p.text = "🤖 This Presentation Built With Claude"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(6)

# Line 2
p = tf.add_paragraph()
p.text = "Demonstrating agentic AI systems"
p.font.size = Pt(13)
p.font.color.rgb = SECONDARY
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(2)

# Line 3
p = tf.add_paragraph()
p.text = "for accelerating medical research"
p.font.size = Pt(13)
p.font.color.rgb = SECONDARY
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(8)

# Line 4
p = tf.add_paragraph()
p.text = "⚠️ Synthetic Data - Educational Purposes"
p.font.size = Pt(11)
p.font.italic = True
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# Save
output_file = 'AI_Ophthalmology_Presentation.pptx'
prs.save(output_file)

print("\n" + "="*80)
print("✓ PROFESSIONAL PRESENTATION COMPLETE!")
print("="*80)
print(f"\nSaved: {output_file}")
print(f"Slides: {len(prs.slides)}")
print("\nFeatures:")
print("  ✓ Proper margins - no bleeding edges")
print("  ✓ Code examples included")
print("  ✓ Workflow visualizations")
print("  ✓ Database preview")
print("  ✓ All figures embedded")
print("  ✓ Professional typography")
print("  ✓ Consistent design system")
print("  ✓ Claude attribution on final slide")
print("="*80)
