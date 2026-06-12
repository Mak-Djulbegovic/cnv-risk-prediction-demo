"""
PowerPoint Presentation Generator - Professional Version
========================================================

Creates a visually appealing PowerPoint presentation for the ophthalmology
residents demonstration on agentic AI systems for research.

⚠️ SYNTHETIC DATA - For Educational Demonstration Only
"""

import subprocess
import sys

# Install python-pptx if not available
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "-q", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

import os

print("=" * 70)
print("GENERATING PROFESSIONAL POWERPOINT PRESENTATION")
print("=" * 70)

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Professional medical theme
COLOR_PRIMARY = RGBColor(41, 128, 185)      # Blue
COLOR_SECONDARY = RGBColor(52, 73, 94)      # Dark blue-grey
COLOR_ACCENT = RGBColor(231, 76, 60)        # Red accent
COLOR_SUCCESS = RGBColor(46, 204, 113)      # Green
COLOR_WARNING = RGBColor(241, 196, 15)      # Yellow
COLOR_LIGHT = RGBColor(236, 240, 241)       # Light grey
COLOR_DARK = RGBColor(44, 62, 80)           # Dark grey
COLOR_WHITE = RGBColor(255, 255, 255)

def add_footer(slide, text="⚠️ SYNTHETIC DATA - Educational Demonstration"):
    """Add footer to slide"""
    left = Inches(0.5)
    top = Inches(7)
    width = Inches(12.333)
    height = Inches(0.4)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(127, 140, 141)
    p.font.italic = True

def add_background_shape(slide):
    """Add decorative background shapes"""
    # Top right accent
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11), Inches(0), Inches(2.5), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    shape.rotation = 45

    # Bottom left accent
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(-0.5), Inches(6), Inches(2), Inches(2)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLOR_LIGHT
    shape2.line.fill.background()
    shape2.rotation = 45

def create_title_slide(prs, title, subtitle):
    """Create professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_PRIMARY
    background.line.fill.background()

    # Eye icon placeholder (using circle)
    eye_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(5.5), Inches(1.2), Inches(2.3), Inches(1.5)
    )
    eye_shape.fill.solid()
    eye_shape.fill.fore_color.rgb = COLOR_WHITE
    eye_shape.line.color.rgb = COLOR_WHITE
    eye_shape.line.width = Pt(3)

    # Pupil
    pupil = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(6.3), Inches(1.6), Inches(0.7), Inches(0.7)
    )
    pupil.fill.solid()
    pupil.fill.fore_color.rgb = COLOR_SECONDARY
    pupil.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Warning banner
    warning_box = slide.shapes.add_textbox(Inches(3), Inches(6.3), Inches(7.333), Inches(0.5))
    warning_frame = warning_box.text_frame
    warning_frame.text = "⚠️ SYNTHETIC DATA - Educational Demonstration Only"
    p = warning_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WARNING
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_section_slide(prs, title, color=COLOR_SECONDARY):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, bullet_points, add_bg=True):
    """Create content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if add_bg:
        add_background_shape(slide)

    # Title background
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLOR_PRIMARY
    title_bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.margin_left = Inches(0.2)
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Content area
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11.333), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]

        # Determine if this is a bullet or sub-bullet
        if point.startswith('  '):
            p.text = point.strip()
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_SECONDARY
        elif point.startswith('    '):
            p.text = point.strip()
            p.level = 2
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_DARK
        elif point == "":
            p.text = ""
        else:
            p.text = point
            p.level = 0
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = COLOR_SECONDARY

        p.space_after = Pt(8)

    add_footer(slide)
    return slide

def create_two_column_slide(prs, title, left_points, right_points):
    """Create slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_shape(slide)

    # Title background
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLOR_PRIMARY
    title_bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5))
    left_frame = left_box.text_frame
    for i, point in enumerate(left_points):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(5.5), Inches(5))
    right_frame = right_box.text_frame
    for i, point in enumerate(right_points):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_SECONDARY
        p.space_after = Pt(6)

    add_footer(slide)
    return slide

def create_image_slide(prs, title, image_path, caption=""):
    """Create slide with large image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLOR_PRIMARY
    title_bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Image
    if os.path.exists(image_path):
        left = Inches(1)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(11.333))

    if caption:
        caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.333), Inches(0.4))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        p = caption_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = COLOR_DARK
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide

print("\nCreating slides...")

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
print("  1. Title slide")
slide1 = create_title_slide(
    prs,
    "AI-Powered Research in Ophthalmology",
    "Agentic Systems in Action: From Data to Clinical Tools\nA Live Demonstration"
)

# ============================================================================
# SLIDE 2: Disclosures
# ============================================================================
print("  2. Disclosures")
slide2 = create_content_slide(
    prs,
    "Disclosures",
    ["No financial disclosures", "", "No conflicts of interest"]
)

# ============================================================================
# SLIDE 3: The AI Landscape
# ============================================================================
print("  3. AI Landscape")
slide3 = create_content_slide(
    prs,
    "The AI Landscape: Many Models Available",
    [
        "Available AI Models:",
        "  • GPT-4 (OpenAI)",
        "  • Gemini (Google)",
        "  • Claude (Anthropic) ← Today's Focus",
        "  • Llama (Meta)",
        "  • And many others...",
        "",
        "Why Claude for Research?",
        "  • Strong reasoning and analytical capabilities",
        "  • Excellent for complex research tasks",
        "  • Agentic capabilities with Claude Code",
        "  • Extended context for large datasets"
    ]
)

# ============================================================================
# SLIDE 4: Evolution
# ============================================================================
print("  4. LLM to Agentic evolution")
slide4 = create_two_column_slide(
    prs,
    "The Evolution: From LLMs to Agentic Systems",
    [
        "Traditional LLMs (2023):",
        "",
        "• Single query → single response",
        "",
        "• No memory between tasks",
        "",
        "• Manual copy-paste workflow",
        "",
        "• You do all the work",
        "",
        "• Like a consultant who just advises"
    ],
    [
        "Agentic Systems (2024-2025):",
        "",
        "• Multi-step autonomous execution",
        "",
        "• Uses tools: code, search, files",
        "",
        "• Remembers context and iterates",
        "",
        "• AI does the implementation",
        "",
        "• Like a research assistant who executes"
    ]
)

# ============================================================================
# SLIDE 5: Capabilities
# ============================================================================
print("  5. Capabilities")
slide5 = create_content_slide(
    prs,
    "What Can Agentic Systems Do for Research?",
    [
        "📊 Data Generation & Simulation",
        "  • Create realistic synthetic datasets",
        "  • Literature-based parameter selection",
        "",
        "📈 Statistical Analysis",
        "  • Complex models (Cox, Random Forest, etc.)",
        "  • Survival analysis, regression, machine learning",
        "",
        "📉 Visualization",
        "  • Publication-quality figures",
        "  • Interactive dashboards",
        "",
        "📚 Literature Integration",
        "  • Search and synthesize research papers",
        "  • Evidence-based methodology",
        "",
        "🌐 Tool Building",
        "  • Interactive web applications",
        "  • Clinical decision support tools"
    ]
)

# ============================================================================
# SLIDE 6: Prompt Engineering
# ============================================================================
print("  6. Prompt engineering")
slide6 = create_content_slide(
    prs,
    "Prompt Engineering: The 4-C Framework",
    [
        "How to Communicate with Agentic Systems:",
        "",
        "1. CONTEXT - Set the stage",
        '   "I\'m studying choroidal rupture outcomes in trauma patients"',
        "",
        "2. CLEAR GOAL - State what you want",
        '   "Create a database of 1000 patients with CNV as the outcome"',
        "",
        "3. CONSTRAINTS - Provide specifications",
        '   "Include 25-30 clinically relevant variables based on literature"',
        "",
        "4. COLLABORATION - Iterate and refine",
        '   "Show me the data structure first, then we\'ll adjust"',
        "",
        "💡 Key Difference: Agentic systems execute autonomously",
        "    You guide the direction → They handle implementation"
    ]
)

# ============================================================================
# SLIDE 7: Today's Demo
# ============================================================================
print("  7. Demo setup")
slide7 = create_content_slide(
    prs,
    "Today's Demo: Predicting CNV After Choroidal Rupture",
    [
        "Clinical Research Question:",
        "  • Which patients develop choroidal neovascularization?",
        "  • How long until CNV develops?",
        "  • Can we build a prediction model?",
        "",
        "What We'll Demonstrate Live:",
        "  1. Generate synthetic database (1000 patients, 33 variables)",
        "  2. Exploratory data analysis",
        "  3. Build predictive models (Cox + Random Forest)",
        "  4. Create publication-quality visualizations",
        "  5. Build interactive web application",
        "",
        "All executed by Claude Code as an autonomous agent",
        "You'll see: Natural language prompts → Complete implementation"
    ]
)

# ============================================================================
# SLIDE 8: Live Demo Section
# ============================================================================
print("  8. Live demo placeholder")
slide8 = create_section_slide(prs, "LIVE DEMONSTRATION", COLOR_ACCENT)

# ============================================================================
# SLIDE 9: Demo Phases
# ============================================================================
print("  9. Demo phases")
slide9 = create_content_slide(
    prs,
    "Demo Phases - Watch the Agent Work",
    [
        "Phase 1: Data Generation (5 min)",
        "  → 1000 patients × 33 variables",
        "  → Show database structure and depth",
        "",
        "Phase 2: Exploratory Analysis (4 min)",
        "  → Descriptive statistics, risk factor identification",
        "",
        "Phase 3: Predictive Modeling (4 min)",
        "  → Cox Proportional Hazards + Random Forest",
        "",
        "Phase 4: Visualizations (2 min)",
        "  → 7 publication-quality figures generated",
        "",
        "Phase 5: Web Application (3-4 min)",
        "  → Interactive clinical tool",
        "  → Live patient case demonstrations"
    ]
)

# ============================================================================
# SLIDE 10: Web App Demo
# ============================================================================
print("  10. Web app features")
slide10 = create_content_slide(
    prs,
    "Interactive Clinical Tool - Live Demo",
    [
        "Patient Risk Calculator Features:",
        "  • Input patient-specific parameters",
        "  • Real-time CNV risk prediction",
        "  • Risk stratification (Low / Medium / High)",
        "  • Contributing factor analysis",
        "  • Clinical recommendations",
        "  • Downloadable patient reports",
        "",
        "We'll Demonstrate Two Cases:",
        "  1. Low-risk patient (~15% CNV probability)",
        "  2. High-risk patient (~75% CNV probability)",
        "",
        "From Research Question → Deployable Tool",
        "All in ONE collaborative session with the agent"
    ]
)

# ============================================================================
# SLIDE 11: Results Overview
# ============================================================================
print("  11. Results overview")
slide11 = create_content_slide(
    prs,
    "Key Findings from Our Analysis",
    [
        "Dataset Generated:",
        "  • 1000 patients, 33 clinical variables",
        "  • CNV developed in 634 patients (63.4%)",
        "  • Mean time to CNV: 6.7 months",
        "",
        "Top Risk Factors Identified:",
        "  1. Foveal involvement (Hazard Ratio = 1.79)",
        "  2. Distance from fovea (closer = higher risk)",
        "  3. Outer retinal disruption severity",
        "  4. Baseline visual acuity",
        "",
        "Model Performance:",
        "  • Cox model C-index: 0.662",
        "  • Random Forest AUC: 0.768",
        "  • Sensitivity: 85% (catches most CNV cases)",
        "  • Enables risk-stratified monitoring protocols"
    ]
)

# ============================================================================
# SLIDE 12-14: Figure Slides
# ============================================================================
print("  12. Kaplan-Meier figure")
if os.path.exists('figure1_kaplan_meier.png'):
    slide12 = create_image_slide(
        prs,
        "Survival Analysis: Time to CNV Development",
        'figure1_kaplan_meier.png',
        "Kaplan-Meier curves stratified by foveal involvement (p < 0.001)"
    )

print("  13. Feature importance figure")
if os.path.exists('figure2_feature_importance.png'):
    slide13 = create_image_slide(
        prs,
        "Top Predictors of CNV Development",
        'figure2_feature_importance.png',
        "Random Forest feature importance analysis"
    )

print("  14. ROC curve figure")
if os.path.exists('figure3_roc_curve.png'):
    slide14 = create_image_slide(
        prs,
        "Model Performance: Predictive Accuracy",
        'figure3_roc_curve.png',
        "ROC curve demonstrating strong discriminative ability (AUC = 0.768)"
    )

# ============================================================================
# SLIDE 15: Code Review
# ============================================================================
print("  15. Code review")
slide15 = create_content_slide(
    prs,
    "Behind the Scenes: Agent-Generated Code",
    [
        "What the Agent Autonomously Created:",
        "",
        "📝 Data Generation Script (~300 lines)",
        "  • Literature-based realistic parameters",
        "  • Complex correlation structures",
        "",
        "📊 Analysis Pipeline (~200 lines)",
        "  • Statistical tests, summary tables",
        "",
        "🤖 Predictive Models (~250 lines)",
        "  • Cox Proportional Hazards implementation",
        "  • Random Forest classifier with cross-validation",
        "",
        "🎨 Visualization Suite (~300 lines)",
        "  • 7 publication-quality figures",
        "",
        "🌐 Web Application (~400 lines)",
        "  • Full interactive Streamlit interface",
        "",
        "YOU didn't write this code - the AGENT did!",
        "All from natural language prompts in the 4-C framework"
    ]
)

# ============================================================================
# SLIDE 16: Other Use Cases
# ============================================================================
print("  16. Other use cases")
slide16 = create_content_slide(
    prs,
    "Beyond This Example: Broader Applications",
    [
        "Research Applications:",
        "  • Systematic reviews & meta-analyses",
        "  • Clinical trial design & sample size calculations",
        "  • Grant writing & literature synthesis",
        "  • Image analysis pipelines (OCT, fundus photos)",
        "  • Multi-center data harmonization",
        "",
        "Clinical Applications:",
        "  • EHR data extraction & cohort identification",
        "  • Custom clinical calculators",
        "  • Outcome prediction tools",
        "  • Quality improvement dashboards",
        "",
        "Education:",
        "  • Interactive case-based learning modules",
        "  • Competency assessment tools",
        "  • Teaching case generation"
    ]
)

# ============================================================================
# SLIDE 17: Practical Tips
# ============================================================================
print("  17. Practical tips")
slide17 = create_content_slide(
    prs,
    "Getting Started: Practical Tips",
    [
        "Best Practices:",
        "  ✓ Use the 4-C framework for clear prompts",
        "  ✓ Let the agent work autonomously - don't micromanage",
        "  ✓ Iterate with feedback - AI learns from your guidance",
        "  ✓ Always verify outputs - AI can make mistakes",
        "",
        "What Agentic Systems ARE Good For:",
        "  • Data exploration and hypothesis generation",
        "  • Rapid prototyping of analyses",
        "  • Tool building and automation",
        "  • Learning new statistical methods",
        "",
        "What They Are NOT:",
        "  ✗ Replacement for biostatistician on final analysis",
        "  ✗ Ready for publication without expert review",
        "  ✗ Substitute for clinical judgment",
        "",
        "Think: Powerful Research Assistant, Not Replacement"
    ]
)

# ============================================================================
# SLIDE 18: Resources
# ============================================================================
print("  18. Resources")
slide18 = create_content_slide(
    prs,
    "Resources & Next Steps",
    [
        "Getting Started Today:",
        "  • Claude.ai - Free tier available for exploration",
        "  • Claude Code CLI - For advanced agentic workflows",
        "  • Cursor / Windsurf - AI-powered development environments",
        "",
        "Important Considerations:",
        "  ⚠️ Review institutional AI use policies",
        "  ⚠️ HIPAA compliance for real patient data",
        "  ⚠️ Data privacy and security protocols",
        "  ⚠️ IRB requirements for AI-assisted research",
        "  ⚠️ Authorship and disclosure guidelines",
        "",
        "Key Takeaway:",
        "Agentic AI systems are available NOW",
        "Start simple → Learn → Scale → Transform your research"
    ]
)

# ============================================================================
# SLIDE 19: Questions
# ============================================================================
print("  19. Questions")
slide19 = create_section_slide(prs, "Questions & Discussion", COLOR_PRIMARY)

# ============================================================================
# SLIDE 20: Closing - Attribution
# ============================================================================
print("  20. Closing with attribution")
slide20 = prs.slides.add_slide(prs.slide_layouts[6])

# Background gradient effect
bg1 = slide20.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(7.5)
)
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_SECONDARY
bg1.line.fill.background()

# Thank you message
thank_you_box = slide20.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You!"
p = thank_you_frame.paragraphs[0]
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide20.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "From Research Question to Clinical Tool in One Session"
p = subtitle_frame.paragraphs[0]
p.font.size = Pt(26)
p.font.color.rgb = COLOR_LIGHT
p.alignment = PP_ALIGN.CENTER

# Attribution box with border
attr_bg = slide20.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(3), Inches(4.5), Inches(7.333), Inches(1.8)
)
attr_bg.fill.solid()
attr_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
attr_bg.line.color.rgb = COLOR_PRIMARY
attr_bg.line.width = Pt(3)

attr_box = slide20.shapes.add_textbox(Inches(3.2), Inches(4.7), Inches(7), Inches(1.5))
attr_frame = attr_box.text_frame
attr_frame.text = "🤖 This Presentation Was Built With Claude\n\n"
attr_frame.text += "Demonstrating the power of agentic AI systems\n"
attr_frame.text += "for accelerating medical research\n\n"
attr_frame.text += "⚠️ All data shown is SYNTHETIC - for educational purposes"

for i, para in enumerate(attr_frame.paragraphs):
    if i == 0:
        para.font.size = Pt(22)
        para.font.bold = True
        para.font.color.rgb = COLOR_PRIMARY
    else:
        para.font.size = Pt(14)
        para.font.color.rgb = COLOR_SECONDARY
    para.alignment = PP_ALIGN.CENTER

# Save presentation
output_file = '/Users/makdjulbegovic/Desktop/Claude_Teach/AI_Ophthalmology_Presentation.pptx'
prs.save(output_file)

print("\n" + "=" * 70)
print(f"✓ PROFESSIONAL PRESENTATION CREATED!")
print("=" * 70)
print(f"\nSaved to: {output_file}")
print(f"Total slides: {len(prs.slides)}")
print("\nPresentation includes:")
print("  • Modern, professional design with color scheme")
print("  • Visual hierarchy and formatting")
print("  • Embedded figures (if available)")
print("  • Two-column layouts for comparisons")
print("  • Section dividers")
print("  • Attribution to Claude at the end")
print("  • Footer on every slide")
print("\n" + "=" * 70)
print("✓ Ready to present with style!")
print("=" * 70)
